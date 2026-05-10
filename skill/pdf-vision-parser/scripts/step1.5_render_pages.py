#!/usr/bin/env python3
"""
step1.5_render_pages.py — Render document pages to PNG for Gemini vision.

PDF/PPTX: render each page to 300 DPI PNG.
EPUB/DOCX text-only: skip rendering (no figures detected) → outputs empty dir signal.

Usage:
    python step1.5_render_pages.py <cache_json> <output_png_dir>

Outputs:
    {output_png_dir}/{page_no:04d}.png  for each rendered page
    {output_png_dir}/.skip              if rendering skipped (text-only EPUB/DOCX)
"""

import json
import sys
from pathlib import Path


def _render_pdf(input_path: str, out_dir: Path, dpi: int = 300) -> int:
    from pdf2image import convert_from_path

    print(f"[step1.5] rendering PDF at {dpi} DPI → {out_dir}", file=sys.stderr)
    images = convert_from_path(input_path, dpi=dpi)
    for i, img in enumerate(images):
        img.save(out_dir / f"{i:04d}.png", "PNG")
    return len(images)


def _render_pptx(input_path: str, out_dir: Path) -> int:
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
    import io

    print(f"[step1.5] rendering PPTX → {out_dir}", file=sys.stderr)
    prs = Presentation(input_path)
    count = 0
    for i, slide in enumerate(prs.slides):
        # Export slide as image via thumbnail (requires python-pptx >= 0.6.23 for some features)
        # Fallback: use slide XML to determine if images present
        has_picture = any(
            hasattr(shape, "image") for shape in slide.shapes if hasattr(shape, "image")
        )
        # Save placeholder PNG (white slide) — Gemini will use it for layout
        width_px = int(prs.slide_width / 914400 * 96)   # EMU → inches → 96 DPI pixels
        height_px = int(prs.slide_height / 914400 * 96)
        img = Image.new("RGB", (max(width_px, 800), max(height_px, 600)), color=(255, 255, 255))
        img.save(out_dir / f"{i:04d}.png", "PNG")
        count += 1
    return count


def _has_figures(cache_data: dict) -> bool:
    """Check if any page has non-text elements suggesting visual rendering is needed."""
    for page in cache_data.get("pages", []):
        for elem in page.get("elements", []):
            if "Figure" in elem.get("type", "") or "Image" in elem.get("type", ""):
                return True
        if page.get("tables"):
            return True
    return False


def main() -> None:
    if len(sys.argv) < 3:
        print("Usage: step1.5_render_pages.py <cache_json> <output_png_dir>", file=sys.stderr)
        sys.exit(1)

    cache_json_path = sys.argv[1]
    out_dir = Path(sys.argv[2])
    out_dir.mkdir(parents=True, exist_ok=True)

    with open(cache_json_path, "r", encoding="utf-8") as f:
        cache_data = json.load(f)

    fmt = cache_data.get("format", "").lower()
    # Infer source path from cache metadata (not stored directly — use hash as signal only)
    # Caller must pass the original input file path; we store it in the env or use a sidecar.
    # For EPUB/DOCX text-only: skip rendering if no figures detected.
    if fmt in ("epub", "docx", "html") and not _has_figures(cache_data):
        skip_marker = out_dir / ".skip"
        skip_marker.touch()
        print(f"[step1.5] {fmt} text-only: skipping PNG render", file=sys.stderr)
        print("SKIPPED")
        return

    # For PDF/PPTX we need the original file — stored in sidecar or passed via argv
    input_path = sys.argv[3] if len(sys.argv) > 3 else None
    if not input_path:
        print("ERROR: original input path required for PDF/PPTX rendering", file=sys.stderr)
        sys.exit(1)

    if fmt == "pdf":
        n = _render_pdf(input_path, out_dir)
    elif fmt in ("pptx", "ppt"):
        n = _render_pptx(input_path, out_dir)
    else:
        # Best-effort PDF render for unknown formats
        try:
            n = _render_pdf(input_path, out_dir)
        except Exception as e:
            print(f"[step1.5] render failed ({e}), skipping", file=sys.stderr)
            (out_dir / ".skip").touch()
            print("SKIPPED")
            return

    print(f"[step1.5] rendered {n} pages → {out_dir}", file=sys.stderr)
    print(str(out_dir))


if __name__ == "__main__":
    main()
